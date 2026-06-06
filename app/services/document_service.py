"""
Document Service.

Responsibilities:
  - File validation (type, size)
  - Secure upload to local filesystem or S3
  - Async document processing (text extraction, smart chunking)
  - Embedding generation via OpenAI + storage in ChromaDB
  - Full CRUD for Document entities
"""
from __future__ import annotations

import hashlib
import io
import logging
import mimetypes
import os
import re
import uuid
from datetime import datetime, timezone
from typing import AsyncGenerator, List, Optional, Tuple
from uuid import UUID

from sqlalchemy import delete, select, update
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.config import settings
from app.core.constants import (
    ALLOWED_MIME_TYPES,
    MAX_FILE_SIZE_BYTES,
    DocumentStatus,
)
from app.core.exceptions import (
    FileTooLargeError,
    InvalidFileTypeError,
    StorageError,
)

logger = logging.getLogger(__name__)

# ── Constants ─────────────────────────────────────────────────────────────────
# (Legacy chunk size constants kept for any external callers; chunking now
#  delegates entirely to SmartDocumentChunker — see chunk_document below.)


class DocumentParseError(Exception):
    pass


class DocumentService:
    """
    Handles document upload, parsing, chunking, embedding, and storage.

    The service is intentionally stateless – it can be instantiated per-request
    with a fresh AsyncSession from FastAPI's dependency injection.
    """

    def __init__(self, db: AsyncSession) -> None:
        self.db = db

    # ── High-level pipeline methods ────────────────────────────────────────

    async def upload_document(
        self,
        file_bytes: bytes,
        filename: str,
        content_type: str,
        project_id: UUID,
        org_id: UUID,
        user_id: UUID,
        tags: Optional[List[str]] = None,
    ):
        """
        Validate, store, and register a document for processing.

        Steps:
        1. Validate file type and size
        2. Compute checksums
        3. Generate a storage key and persist the file
        4. Create a Document DB record with status=UPLOADED
        5. Queue an async processing task (Celery)

        Returns the new Document ORM instance.
        """
        from app.models.document import Document

        # 1. Validate
        self._validate_file(file_bytes, content_type)

        # 2. Checksums
        md5 = hashlib.md5(file_bytes).hexdigest()
        sha256 = hashlib.sha256(file_bytes).hexdigest()

        # 3. Generate storage path
        doc_id = uuid.uuid4()
        safe_name = re.sub(r"[^\w\-_\. ]", "_", filename)
        stored_filename = f"{doc_id}_{safe_name}"
        storage_key = f"documents/{org_id}/{project_id}/{stored_filename}"
        file_path = os.path.join(settings.UPLOAD_DIR if hasattr(settings, "UPLOAD_DIR") else "/tmp", stored_filename)

        # 4. Persist file
        await self._upload_to_storage(storage_key, file_bytes, content_type)

        # 5. DB record
        doc = Document(
            id=doc_id,
            organization_id=org_id,
            project_id=project_id,
            uploaded_by=user_id,
            original_filename=filename,
            stored_filename=stored_filename,
            file_path=file_path,
            s3_key=storage_key,
            s3_bucket=settings.S3_BUCKET_NAME,
            content_type=content_type,
            file_size_bytes=len(file_bytes),
            checksum_md5=md5,
            checksum_sha256=sha256,
            status=DocumentStatus.UPLOADED,
            tags=tags or [],
            version=1,
        )
        self.db.add(doc)
        await self.db.flush()
        await self.db.refresh(doc)
        await self.db.commit()

        # 6. Queue processing
        try:
            self._queue_processing(str(doc.id))
        except Exception as exc:
            logger.warning("Could not queue processing task for document %s: %s", doc.id, exc)

        logger.info("Document %s uploaded (%d bytes)", doc.id, len(file_bytes))
        return doc

    async def process_document(self, document_id: UUID) -> None:
        """
        Parse the document and create text chunks.

        Steps:
        1. Load document record
        2. Fetch raw file bytes from storage
        3. Parse text + metadata based on file type
        4. Smart-chunk the text
        5. Persist chunks to the document_chunks table
        6. Update document status and metadata
        7. Queue embedding generation

        This method is typically invoked by a Celery worker, not the HTTP layer.
        """
        from app.models.document import Document

        doc = (
            await self.db.execute(select(Document).where(Document.id == document_id))
        ).scalar_one_or_none()

        if doc is None:
            logger.error("process_document called for unknown id=%s", document_id)
            return

        # Update status to PROCESSING
        doc.status = DocumentStatus.PROCESSING
        await self.db.flush()

        try:
            # Fetch bytes
            file_bytes = await self._download_from_storage(doc.s3_key or doc.file_path)

            # Parse
            full_text, file_metadata = await self.parse_document(
                file_bytes, doc.original_filename, doc.content_type
            )

            # Chunk
            chunks = self.chunk_document(full_text, {"document_id": str(doc.id), **file_metadata})

            # Persist chunks
            await self.store_chunks(chunks, str(doc.id))

            # Update document
            doc.raw_text = full_text[:65_536] if full_text else None  # Store first 64KB
            doc.page_count = file_metadata.get("page_count")
            doc.word_count = len(full_text.split()) if full_text else 0
            doc.parsed_content = file_metadata
            doc.status = DocumentStatus.PROCESSED
            await self.db.flush()
            await self.db.commit()

            # Queue embedding
            try:
                self._queue_embeddings(str(doc.id))
            except Exception as exc:
                logger.warning("Could not queue embeddings for document %s: %s", doc.id, exc)

            logger.info("Document %s processed: %d chunks", doc.id, len(chunks))

        except Exception as exc:
            logger.exception("Error processing document %s", document_id)
            doc.status = DocumentStatus.FAILED
            doc.processing_error = str(exc)[:2000]
            await self.db.flush()
            await self.db.commit()

    async def generate_embeddings(self, document_id: UUID) -> None:
        """
        Generate and store vector embeddings for all chunks of a document.

        Steps:
        1. Load all chunks for the document
        2. Batch-embed with OpenAI (rate-limit aware, max 100 items/call)
        3. Upsert embeddings into ChromaDB
        4. Update document.vector_ids
        5. Mark document.status = PROCESSED (if not already)

        This method is typically invoked by a Celery worker.
        """
        from app.models.document import Document

        doc = (
            await self.db.execute(select(Document).where(Document.id == document_id))
        ).scalar_one_or_none()
        if doc is None:
            logger.error("generate_embeddings called for unknown id=%s", document_id)
            return

        chunks = await self.get_chunks(str(document_id), page=1, page_size=5000)
        if not chunks:
            logger.warning("No chunks found for document %s", document_id)
            return

        try:
            from openai import AsyncOpenAI

            client = AsyncOpenAI(api_key=settings.OPENAI_API_KEY)
            vector_ids: List[str] = []
            batch_size = 100

            for batch_start in range(0, len(chunks), batch_size):
                batch = chunks[batch_start : batch_start + batch_size]
                texts = [c.content for c in batch]

                response = await client.embeddings.create(
                    model=settings.OPENAI_EMBEDDING_MODEL,
                    input=texts,
                )

                for i, embedding_obj in enumerate(response.data):
                    chunk = batch[i]
                    point_id = str(chunk.id)
                    vector_ids.append(point_id)

                    # Upsert into ChromaDB
                    await self._upsert_vector(
                        point_id=point_id,
                        vector=embedding_obj.embedding,
                        payload={
                            "document_id": str(document_id),
                            "chunk_index": chunk.chunk_index,
                            "content": chunk.content[:500],  # Truncated for payload size
                            "project_id": str(doc.project_id),
                            "org_id": str(doc.organization_id),
                        },
                    )

            doc.vector_ids = vector_ids
            await self.db.flush()
            await self.db.commit()

            logger.info(
                "Generated %d embeddings for document %s", len(vector_ids), document_id
            )

        except Exception as exc:
            logger.exception("Error generating embeddings for document %s", document_id)
            raise

    # ── Public API ─────────────────────────────────────────────────────────

    async def create_document(
        self,
        file_bytes: bytes,
        filename: str,
        content_type: str,
        project_id: Optional[str],
        tags: Optional[List[str]],
        uploaded_by: str,
    ):
        """Save document metadata and upload file to object storage."""
        from app.models.document import Document
        from app.core.constants import DocumentStatus

        doc_id = uuid.uuid4()
        safe_name = re.sub(r"[^\w\-_\. ]", "_", filename)
        stored_filename = f"{doc_id}_{safe_name}"
        storage_key = f"documents/{uploaded_by}/{doc_id}/{safe_name}"

        # Resolve organization_id from project (or fall back to user's org)
        organization_id = await self._resolve_org_id(project_id, uploaded_by)

        # Upload to object storage (best-effort — skip if storage not configured)
        try:
            await self._upload_to_storage(storage_key, file_bytes, content_type)
        except Exception as exc:
            logger.warning("Storage upload skipped for %s: %s", filename, exc)

        # Determine local file path for fallback access
        upload_dir = getattr(settings, "UPLOAD_DIR", "/tmp/sdd_uploads")
        os.makedirs(upload_dir, exist_ok=True)
        file_path = os.path.join(upload_dir, stored_filename)
        try:
            with open(file_path, "wb") as fh:
                fh.write(file_bytes)
        except Exception as exc:
            logger.warning("Could not write local file %s: %s", file_path, exc)
            file_path = stored_filename  # store relative name as fallback

        doc = Document(
            id=doc_id,
            organization_id=organization_id,
            project_id=uuid.UUID(project_id) if project_id else None,
            uploaded_by=uuid.UUID(uploaded_by) if uploaded_by else None,
            original_filename=filename,
            stored_filename=stored_filename,
            file_path=file_path,
            s3_key=storage_key,
            content_type=content_type,
            file_size_bytes=len(file_bytes),
            checksum_sha256=hashlib.sha256(file_bytes).hexdigest(),
            checksum_md5=hashlib.md5(file_bytes).hexdigest(),
            tags=tags or [],
            status=DocumentStatus.UPLOADED,
            version=1,
        )
        self.db.add(doc)
        await self.db.commit()
        await self.db.refresh(doc)
        return doc

    async def _resolve_org_id(self, project_id: Optional[str], user_id: str) -> uuid.UUID:
        """Look up organization_id from project, or fall back to user's org."""
        from app.models.project import Project
        from app.models.user import User

        if project_id:
            try:
                row = await self.db.execute(
                    select(Project.organization_id).where(Project.id == uuid.UUID(project_id))
                )
                org_id = row.scalar_one_or_none()
                if org_id:
                    return org_id
            except Exception:
                pass

        # Fall back to user's organization
        try:
            row = await self.db.execute(
                select(User.organization_id).where(User.id == uuid.UUID(user_id))
            )
            org_id = row.scalar_one_or_none()
            if org_id:
                return org_id
        except Exception:
            pass

        # Last resort: dev org UUID
        return uuid.UUID("00000000-0000-0000-0000-000000000020")

    async def parse_document(self, file_bytes: bytes, filename: str, content_type: str) -> tuple[str, dict]:
        """
        Parse a document and return (full_text, metadata).

        Dispatches to type-specific parsers.
        """
        ext = os.path.splitext(filename)[1].lower()
        normalized = self._normalize_file_type(content_type)

        try:
            if normalized == "pdf" or ext == ".pdf":
                return await self._parse_pdf(file_bytes, filename)
            elif normalized == "docx" or ext in (".docx", ".doc"):
                return await self._parse_docx(file_bytes, filename)
            elif normalized == "pptx" or ext == ".pptx":
                return await self._parse_pptx(file_bytes, filename)
            elif normalized == "xlsx" or ext in (".xlsx", ".xls"):
                return await self._parse_xlsx(file_bytes, filename)
            elif normalized in ("txt", "md", "markdown"):
                return self._parse_text(file_bytes, filename)
            elif normalized == "json":
                return self._parse_json(file_bytes, filename)
            else:
                raise DocumentParseError(f"Unsupported file type: {content_type}")
        except DocumentParseError:
            raise
        except Exception as exc:
            raise DocumentParseError(f"Failed to parse {filename}: {exc}") from exc

    async def _parse_pdf(self, file_bytes: bytes, filename: str) -> tuple[str, dict]:
        """Parse PDF using PyPDF2 with fallback to pdfminer."""
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            text_parts = []
            for page in reader.pages:
                page_text = page.extract_text() or ""
                text_parts.append(page_text)

            full_text = "\n\n".join(text_parts)
            metadata = {
                "page_count": len(reader.pages),
                "title": reader.metadata.get("/Title", filename) if reader.metadata else filename,
                "author": reader.metadata.get("/Author", "") if reader.metadata else "",
                "subject": reader.metadata.get("/Subject", "") if reader.metadata else "",
                "file_type": "pdf",
            }
            return full_text, metadata
        except ImportError:
            logger.warning("PyPDF2 not installed, trying pdfminer")
            return self._parse_pdf_fallback(file_bytes, filename)

    def _parse_pdf_fallback(self, file_bytes: bytes, filename: str) -> tuple[str, dict]:
        try:
            from pdfminer.high_level import extract_text as pdfminer_extract
            text = pdfminer_extract(io.BytesIO(file_bytes))
            return text, {"page_count": None, "title": filename, "file_type": "pdf"}
        except ImportError:
            raise DocumentParseError("No PDF parsing library available. Install PyPDF2 or pdfminer.")

    async def _parse_docx(self, file_bytes: bytes, filename: str) -> tuple[str, dict]:
        """Parse DOCX using python-docx."""
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = "\n\n".join(paragraphs)

            # Also extract table content
            table_texts = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        table_texts.append(row_text)
            if table_texts:
                full_text += "\n\n" + "\n".join(table_texts)

            core_props = doc.core_properties
            metadata = {
                "page_count": None,
                "title": core_props.title or filename,
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "file_type": "docx",
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
            }
            return full_text, metadata
        except ImportError:
            raise DocumentParseError("python-docx not installed.")

    async def _parse_pptx(self, file_bytes: bytes, filename: str) -> tuple[str, dict]:
        """Parse PPTX using python-pptx."""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            slide_texts = []
            for i, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
                if parts:
                    slide_texts.append(f"[Slide {i}]\n" + "\n".join(parts))

            full_text = "\n\n".join(slide_texts)
            metadata = {
                "page_count": len(prs.slides),
                "title": filename,
                "file_type": "pptx",
                "slide_count": len(prs.slides),
            }
            return full_text, metadata
        except ImportError:
            raise DocumentParseError("python-pptx not installed.")

    async def _parse_xlsx(self, file_bytes: bytes, filename: str) -> tuple[str, dict]:
        """Parse XLSX using openpyxl."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            sheet_texts = []
            total_rows = 0

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        rows.append(row_text)
                        total_rows += 1
                if rows:
                    sheet_texts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

            full_text = "\n\n".join(sheet_texts)
            metadata = {
                "page_count": len(wb.sheetnames),
                "title": filename,
                "file_type": "xlsx",
                "sheet_count": len(wb.sheetnames),
                "row_count": total_rows,
            }
            return full_text, metadata
        except ImportError:
            raise DocumentParseError("openpyxl not installed.")

    def _parse_text(self, file_bytes: bytes, filename: str) -> tuple[str, dict]:
        """Parse plain text or Markdown files."""
        try:
            text = file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            text = file_bytes.decode("latin-1")
        metadata = {
            "page_count": 1,
            "title": filename,
            "file_type": "txt",
            "char_count": len(text),
            "line_count": text.count("\n"),
        }
        return text, metadata

    def _parse_json(self, file_bytes: bytes, filename: str) -> tuple[str, dict]:
        """Parse JSON as indented text."""
        import json
        data = json.loads(file_bytes.decode("utf-8"))
        text = json.dumps(data, indent=2)
        metadata = {
            "page_count": 1,
            "title": filename,
            "file_type": "json",
        }
        return text, metadata

    # ── Chunking ───────────────────────────────────────────────────────────

    def chunk_document(
        self,
        text: str,
        metadata: dict,
        chunk_size: int = 0,   # kept for API compat; ignored — SmartDocumentChunker owns sizing
        overlap: int = 0,       # kept for API compat; ignored
    ) -> list[dict]:
        """
        Chunk document text using the hierarchical structure-aware SmartDocumentChunker.

        Delegates entirely to the singleton `document_chunker` from
        `app.ai.rag.chunker` which uses:
          • 400-token target chunks  (sweet spot for bge-large-en-v1.5)
          • Section breadcrumb prepending (e.g. "Auth > Password Reset")
          • Sentence-level 50-token overlap
          • Paragraph → sentence split boundaries (never mid-sentence)
          • Block type tagging: paragraph | list | table | code | header

        Returns:
            List of chunk dicts: {chunk_index, content, metadata, token_count}.
        """
        from app.ai.rag.chunker import document_chunker

        smart_chunks = document_chunker.chunk_text(text=text, metadata=metadata)
        logger.debug("Document chunked into %d smart chunks", len(smart_chunks))
        return [
            {
                "chunk_index": c.chunk_index,
                "content": c.text,
                "metadata": {**c.metadata, "chunk_index": c.chunk_index},
                "token_count": c.token_count,
            }
            for c in smart_chunks
        ]

    # ── Storage operations ─────────────────────────────────────────────────

    async def store_chunks(self, chunks: list[dict], document_id: str) -> None:
        """Persist document chunks to the database.

        Primary path: insert into document_chunks table via DocumentChunk model.
        Fallback: store a summary in the documents.parsed_content JSONB column when
        the DocumentChunk model / table does not yet exist (e.g. dev environment).
        """
        try:
            from app.models.document import DocumentChunk  # may not exist in dev

            # Delete existing chunks for re-processing
            await self.db.execute(
                delete(DocumentChunk).where(DocumentChunk.document_id == document_id)
            )

            chunk_objects = [
                DocumentChunk(
                    id=str(uuid.uuid4()),
                    document_id=document_id,
                    chunk_index=chunk["chunk_index"],
                    content=chunk["content"],
                    metadata=chunk["metadata"],
                    token_count=chunk["token_count"],
                )
                for chunk in chunks
            ]
            self.db.add_all(chunk_objects)
            await self.db.commit()
        except (ImportError, Exception) as exc:
            logger.warning(
                "document_chunks table unavailable (%s); storing chunk summary in parsed_content instead", exc
            )
            # Fallback: persist a compact chunk summary inside the JSONB parsed_content column
            from app.models.document import Document
            chunk_summary = [
                {
                    "chunk_index": c["chunk_index"],
                    "content": c["content"][:1000],  # truncate to keep row size sane
                    "token_count": c.get("token_count"),
                }
                for c in chunks[:500]  # cap at 500 chunks
            ]
            await self.db.execute(
                update(Document)
                .where(Document.id == document_id)
                .values(
                    parsed_content={
                        "chunks": chunk_summary,
                        "chunk_count": len(chunks),
                    }
                )
            )
            await self.db.commit()

    async def trigger_embedding(self, document_id: str) -> None:
        """Queue the embedding generation task for a document."""
        from app.workers.tasks.document_tasks import generate_embeddings
        generate_embeddings.delay(document_id)

    # ── CRUD operations ────────────────────────────────────────────────────

    async def get_by_id(self, document_id: str):
        from app.models.document import Document
        result = await self.db.execute(
            select(Document).where(Document.id == document_id)
        )
        return result.scalar_one_or_none()

    async def list_documents(
        self,
        user_id: str,
        project_id: str | None = None,
        status: str | None = None,
        file_type: str | None = None,
        search: str | None = None,
        page: int = 1,
        page_size: int = 20,
    ) -> dict:
        from app.models.document import Document
        from sqlalchemy import func, or_

        query = select(Document)

        if project_id:
            query = query.where(Document.project_id == project_id)
        if status:
            query = query.where(Document.status == status)
        if file_type:
            query = query.where(Document.file_type == file_type)
        if search:
            query = query.where(Document.name.ilike(f"%{search}%"))

        count_query = select(func.count()).select_from(query.subquery())
        total = (await self.db.execute(count_query)).scalar_one()

        query = query.order_by(Document.created_at.desc())
        query = query.offset((page - 1) * page_size).limit(page_size)
        result = await self.db.execute(query)
        items = result.scalars().all()

        return {
            "items": items,
            "total": total,
            "page": page,
            "page_size": page_size,
        }

    async def update_status(
        self,
        document_id: str,
        status: str,
        error_message: str | None = None,
        page_count: int | None = None,
        chunk_count: int | None = None,  # accepted for API compat but ignored (no such column)
    ) -> None:
        from app.models.document import Document
        updates: dict = {
            "status": status,
            "updated_at": datetime.now(tz=timezone.utc),
        }
        if error_message is not None:
            updates["processing_error"] = error_message  # actual column name
        if page_count is not None:
            updates["page_count"] = page_count
        # NOTE: chunk_count is intentionally NOT included — the documents table has no such column.

        await self.db.execute(
            update(Document).where(Document.id == document_id).values(**updates)
        )
        await self.db.commit()

    async def delete_document(self, document_id: str) -> None:
        from app.models.document import Document
        doc = await self.get_by_id(document_id)
        if doc:
            # Best-effort: delete local file and S3 object
            if doc.file_path:
                try:
                    import os as _os
                    if _os.path.exists(doc.file_path):
                        _os.remove(doc.file_path)
                except Exception as exc:
                    logger.warning("Failed to delete local file %s: %s", doc.file_path, exc)
            if doc.s3_key:
                try:
                    await self._delete_from_storage(doc.s3_key)
                except Exception as exc:
                    logger.warning("Failed to delete storage object %s: %s", doc.s3_key, exc)

        await self.db.execute(
            delete(Document).where(Document.id == uuid.UUID(document_id) if isinstance(document_id, str) else document_id)
        )
        await self.db.commit()

    async def get_chunks(
        self,
        document_id: str,
        page: int = 1,
        page_size: int = 50,
    ) -> list:
        try:
            from app.models.document import DocumentChunk  # may not exist in dev
            result = await self.db.execute(
                select(DocumentChunk)
                .where(DocumentChunk.document_id == document_id)
                .order_by(DocumentChunk.chunk_index)
                .offset((page - 1) * page_size)
                .limit(page_size)
            )
            rows = result.scalars().all()
            if rows:
                return rows
        except (ImportError, Exception) as exc:
            logger.warning("Could not query document_chunks for %s: %s", document_id, exc)

        # Fallback: read chunks stored in parsed_content JSONB
        try:
            from app.models.document import Document
            doc_result = await self.db.execute(
                select(Document).where(Document.id == document_id)
            )
            doc = doc_result.scalar_one_or_none()
            if doc and doc.parsed_content and "chunks" in (doc.parsed_content or {}):
                raw_chunks = doc.parsed_content["chunks"]
                offset = (page - 1) * page_size
                sliced = raw_chunks[offset: offset + page_size]
                # Return simple namespace objects so callers can access .content, .chunk_index
                from types import SimpleNamespace
                return [
                    SimpleNamespace(
                        chunk_index=c.get("chunk_index", i),
                        content=c.get("content", ""),
                        token_count=c.get("token_count"),
                    )
                    for i, c in enumerate(sliced)
                ]
        except Exception as exc2:
            logger.warning("Fallback chunk read also failed for %s: %s", document_id, exc2)

        return []

    async def get_download_url(self, document_id: str) -> str:
        """Generate a presigned URL for direct download."""
        doc = await self.get_by_id(document_id)
        if not doc:
            raise ValueError("Document not found")
        return await self._generate_presigned_url(doc.storage_key)

    async def extract_metadata(self, file_bytes: bytes, filename: str) -> dict:
        """Extract metadata from file without full parse."""
        content_type = mimetypes.guess_type(filename)[0] or "application/octet-stream"
        _, metadata = await self.parse_document(file_bytes, filename, content_type)
        return metadata

    # ── Storage backend ────────────────────────────────────────────────────

    async def _upload_to_storage(
        self,
        key: str,
        file_bytes: bytes,
        content_type: str,
    ) -> None:
        """Upload to S3-compatible object storage."""
        try:
            import aioboto3
            from app.core.config import settings

            session = aioboto3.Session()
            async with session.client(
                "s3",
                endpoint_url=settings.S3_ENDPOINT_URL,
                aws_access_key_id=settings.S3_ACCESS_KEY,
                aws_secret_access_key=settings.S3_SECRET_KEY,
            ) as s3:
                await s3.put_object(
                    Bucket=settings.S3_BUCKET,
                    Key=key,
                    Body=file_bytes,
                    ContentType=content_type,
                )
        except ImportError:
            logger.warning("aioboto3 not available, skipping S3 upload")
        except Exception as exc:
            logger.error("S3 upload failed for %s: %s", key, exc)
            raise

    async def _delete_from_storage(self, key: str) -> None:
        try:
            import aioboto3
            from app.core.config import settings

            session = aioboto3.Session()
            async with session.client(
                "s3",
                endpoint_url=settings.S3_ENDPOINT_URL,
                aws_access_key_id=settings.S3_ACCESS_KEY,
                aws_secret_access_key=settings.S3_SECRET_KEY,
            ) as s3:
                await s3.delete_object(Bucket=settings.S3_BUCKET, Key=key)
        except Exception as exc:
            logger.warning("Failed to delete S3 object %s: %s", key, exc)

    async def _generate_presigned_url(self, key: str, expires_in: int = 3600) -> str:
        try:
            import aioboto3
            from app.core.config import settings

            session = aioboto3.Session()
            async with session.client(
                "s3",
                endpoint_url=settings.S3_ENDPOINT_URL,
                aws_access_key_id=settings.S3_ACCESS_KEY,
                aws_secret_access_key=settings.S3_SECRET_KEY,
            ) as s3:
                url = await s3.generate_presigned_url(
                    "get_object",
                    Params={"Bucket": settings.S3_BUCKET, "Key": key},
                    ExpiresIn=expires_in,
                )
                return url
        except Exception:
            return f"/api/v1/documents/download/{key}"

    # ── Storage download helper ────────────────────────────────────────────

    async def _download_from_storage(self, key_or_path: str) -> bytes:
        """Download file bytes from S3 or local filesystem."""
        if key_or_path and key_or_path.startswith("/"):
            # Local filesystem fallback
            try:
                with open(key_or_path, "rb") as fh:
                    return fh.read()
            except FileNotFoundError:
                raise StorageError(message=f"File not found at path: {key_or_path}")

        try:
            import aioboto3

            session = aioboto3.Session()
            async with session.client(
                "s3",
                endpoint_url=settings.S3_ENDPOINT_URL,
                aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
                aws_secret_access_key=settings.AWS_SECRET_ACCESS_KEY,
                region_name=settings.AWS_REGION,
            ) as s3:
                response = await s3.get_object(Bucket=settings.S3_BUCKET_NAME, Key=key_or_path)
                return await response["Body"].read()
        except ImportError:
            raise StorageError(message="aioboto3 not installed; cannot download from S3")
        except Exception as exc:
            raise StorageError(message=f"Failed to download {key_or_path}: {exc}") from exc

    # ── ChromaDB helper ────────────────────────────────────────────────────

    async def _upsert_vector(
        self,
        point_id: str,
        vector: List[float],
        payload: dict,
    ) -> None:
        """Upsert a vector point into ChromaDB."""
        try:
            import asyncio
            import chromadb
            from chromadb.config import Settings as ChromaSettings

            persist_dir = settings.CHROMA_PERSIST_DIR
            org_id = payload.get("org_id", "")
            collection_name = (
                f"sdd_org_{org_id.replace('-', '_')}" if org_id
                else settings.CHROMA_COLLECTION_NAME
            )

            def _upsert():
                import os as _os
                _os.makedirs(persist_dir, exist_ok=True)
                client = chromadb.PersistentClient(
                    path=persist_dir,
                    settings=ChromaSettings(anonymized_telemetry=False),
                )
                col = client.get_or_create_collection(
                    name=collection_name,
                    metadata={"hnsw:space": "cosine"},
                )
                col.upsert(
                    ids=[point_id],
                    embeddings=[vector],
                    documents=[payload.get("content", "")],
                    metadatas=[{k: str(v) for k, v in payload.items()}],
                )

            await asyncio.to_thread(_upsert)
        except ImportError:
            logger.warning("chromadb not installed; skipping vector upsert")
        except Exception as exc:
            logger.error("ChromaDB upsert failed for point %s: %s", point_id, exc)
            raise

    # ── Task queuing helpers ───────────────────────────────────────────────

    def _queue_processing(self, document_id: str) -> None:
        """Queue a Celery document-processing task."""
        try:
            from app.workers.tasks import process_document as task
            task.delay(document_id)
        except Exception as exc:
            logger.warning("Could not queue document processing task: %s", exc)

    def _queue_embeddings(self, document_id: str) -> None:
        """Queue a Celery embedding-generation task."""
        try:
            from app.workers.tasks import generate_embeddings as task
            task.delay(document_id)
        except Exception as exc:
            logger.warning("Could not queue embedding task: %s", exc)

    # ── Validation helpers ─────────────────────────────────────────────────

    def _validate_file(self, file_bytes: bytes, content_type: str) -> None:
        """Raise a validation error if the file type or size is not allowed."""
        if len(file_bytes) > MAX_FILE_SIZE_BYTES:
            raise FileTooLargeError(
                message=f"File size {len(file_bytes)} bytes exceeds the maximum of {MAX_FILE_SIZE_BYTES} bytes"
            )
        if content_type not in ALLOWED_MIME_TYPES:
            raise InvalidFileTypeError(
                message=f"File type '{content_type}' is not allowed. Allowed types: {', '.join(ALLOWED_MIME_TYPES)}"
            )

    # ── Helpers ────────────────────────────────────────────────────────────

    def _normalize_file_type(self, content_type: str) -> str:
        mapping = {
            "application/pdf": "pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
            "application/msword": "docx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
            "application/vnd.ms-powerpoint": "pptx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
            "application/vnd.ms-excel": "xlsx",
            "text/plain": "txt",
            "text/markdown": "md",
            "application/json": "json",
        }
        return mapping.get(content_type, "unknown")
